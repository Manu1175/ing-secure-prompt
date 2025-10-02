# Creates SecurePrompt_ING_M1_FINAL.pptx in the current folder.
# If IMAGE paths are missing, it still builds (places a small "(missing image)" label).

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from PIL import Image
from datetime import datetime
import os

# ---- (OPTIONAL) set these to your local screenshots; leave as-is if none ----
IMAGE1 = "Interactive_Scrubbing.png"    # or an absolute path
IMAGE2 = "Audit_Dashboard.png"
IMAGE3 = "API_Docs_Swagger.png"

OUTFILE = "SecurePrompt_ING_M1_FINAL.pptx"

ING_ORANGE = RGBColor(255, 98, 0)
TEXT = RGBColor(25, 25, 25)
BORDER = RGBColor(180, 180, 180)
PANEL = RGBColor(255, 255, 255)
FRAME = RGBColor(248, 248, 248)

SAFE_L, SAFE_T, SAFE_W, SAFE_H = 0.6, 1.35, 9.8, 4.9  # inches
prs = Presentation()

def set_title(slide, text, size=42):
    slide.shapes.title.text = text
    f = slide.shapes.title.text_frame.paragraphs[0].font
    f.size = Pt(size); f.color.rgb = ING_ORANGE

def title_only(text):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(s, text, 36)
    return s

def notes(slide, txt):
    n = slide.notes_slide.notes_text_frame; n.clear(); n.text = txt

def panel(slide, l, t, w, h):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = FRAME
    shp.line.color.rgb = BORDER
    return shp

def bullets(slide, l, t, w, h, title, items, bullet_size=18):
    shp = panel(slide, l, t, w, h)
    tf = shp.text_frame; tf.clear()
    tf.text = title
    p = tf.paragraphs[0]; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = TEXT
    for it in items:
        q = tf.add_paragraph(); q.text = "• " + it; q.font.size = Pt(bullet_size); q.font.color.rgb = TEXT
    return shp

def box(slide, l, t, w, h, title, body=None):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = PANEL; shp.line.color.rgb = BORDER
    tf = shp.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.bold = True; p.font.size = Pt(14); p.font.color.rgb = TEXT
    if body:
        r = tf.add_paragraph(); r.text = body; r.font.size = Pt(12); r.font.color.rgb = TEXT
    return shp

def connect(slide, a, b):
    con = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, a.left + a.width, a.top + a.height/2, b.left, b.top + b.height/2)
    con.line.color.rgb = BORDER

def image_fit(slide, path, l, t, w, h, caption=None):
    fr = panel(slide, l, t, w, h)
    if os.path.exists(path):
        img = Image.open(path); iw, ih = img.size
        max_w, max_h = (w - 0.2), (h - 0.3)
        scale = min((max_w*96)/iw, (max_h*96)/ih)
        w_in, h_in = iw/96*scale, ih/96*scale
        slide.shapes.add_picture(path, Inches(l + (w - w_in)/2), Inches(t + (h - h_in)/2 - 0.05), width=Inches(w_in), height=Inches(h_in))
    else:
        fr.text_frame.text = f"(missing) {path}"
    if caption:
        cap = slide.shapes.add_textbox(Inches(l), Inches(t + h - 0.28), Inches(w), Inches(0.28)).text_frame
        cap.text = caption; cap.paragraphs[0].font.size = Pt(12); cap.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

# -------- Slides --------
s = prs.slides.add_slide(prs.slide_layouts[0])
set_title(s, "SecurePrompt (ING) — Milestone 1", 42)
s.placeholders[1].text = f"Sanitization architecture • workflows • KPI • {datetime.now():%d %b %Y}"
f = s.placeholders[1].text_frame.paragraphs[0].font; f.size = Pt(20); f.color.rgb = TEXT
notes(s, "We’ll show explainable scrubbing, role-gated de-scrub, and append-only audit; PNG screenshots are priority.")

s = title_only("Executive summary (from kickoff)")
bullets(s, SAFE_L, SAFE_T, SAFE_W, SAFE_H, "What we are building", [
    "Production-grade prompt & file scrubber for banking LLMs",
    "Explainable detections with confidence; replace originals with identifiers",
    "Adaptive sensitivity (C2–C4) and role-gated de-scrub with justification",
    "Append-only audit with receipts; highlight malicious prompts",
    "File scope: PDF, DOCX, TXT, HTML; priority: screenshots (PNG)",
], bullet_size=18)
notes(s, "Directly from kickoff.")

s = title_only("Architecture overview + Policy (C2–C4)")
bw, bh, gap = 1.8, 0.95, 0.3
x, y = SAFE_L, SAFE_T + 0.15
clients = box(s, x, y, bw, bh, "Clients", "UI • Swagger • SDK"); x += bw + gap
api = box(s, x, y, bw, bh, "FastAPI", "/scrub /files /descrub /audit"); x += bw + gap
scr = box(s, x, y, bw, bh, "Scrub", "detect → explain → replace"); x += bw + gap
pol = box(s, x, y, bw, bh, "Policy", "C2/C3/C4 matrix"); x += bw + gap
aud = box(s, x, y, bw, bh, "Audit", "hash chain • receipts")
ds = box(s, SAFE_L + bw + gap + bw/2, y + bh + 0.5, bw, bh, "De-scrub", "role + justification")
rc = box(s, SAFE_L + (bw+gap)*4, y + bh + 0.5, bw, bh, "Receipts", "JSON/JSONL")
for a,b in [(clients,api),(api,scr),(scr,pol),(pol,aud)]: connect(s,a,b)
s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ds.left + ds.width/2, ds.top, pol.left + pol.width/2, pol.top + pol.height).line.color.rgb = BORDER
s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, aud.left + aud.width/2, aud.top + aud.height, rc.left + rc.width/2, rc.top).line.color.rgb = BORDER
bullets(s, SAFE_L + SAFE_W/2 + 0.1, SAFE_T + 2.2, SAFE_W/2 - 0.1, SAFE_H - 2.3, "Policy usage (examples)", [
    "Map labels → clearance: PAN/IBAN=C4; address/email/phone=C3; name=C2 (confirm with ING)",
    "Runtime: mask entities above user’s clearance",
    "Per-label toggles for temporary tuning during KPI runs",
], bullet_size=16)

s = title_only("Workflow — text scrubbing")
bxw, bxh, g = 2.1, 0.95, 0.35
x, y = SAFE_L, SAFE_T + 0.2
i1 = box(s, x, y, bxw, bxh, "Input", "prompt/response"); x += bxw + g
i2 = box(s, x, y, bxw, bxh, "Detectors", "regex + validators"); x += bxw + g
i3 = box(s, x, y, bxw, bxh, "Explain & score", "label • rule • confidence • C-level"); x += bxw + g
i4 = box(s, x, y, bxw, bxh, "Replace", "C{lvl}::LABEL::hash"); x += bxw + g
i5 = box(s, x, y, bxw, bxh, "Output", "sanitized text")
for a,b in [(i1,i2),(i2,i3),(i3,i4),(i4,i5)]: connect(s,a,b)
bullets(s, SAFE_L, SAFE_T + 2.2, SAFE_W, SAFE_H - 2.35, "Highlights", [
    "Rules-first for M1 (reliable & fast); ML assist post-M1",
    "Human explanation per hit (e.g., 'detected as IBAN')",
    "Confidence per hit feeds KPI and operator review",
], bullet_size=16)

s = title_only("Workflow — file & OCR")
bw2, bh2, g2 = 2.2, 0.95, 0.3
x, y = SAFE_L, SAFE_T + 0.15
f1 = box(s, x, y, bw2, bh2, "Sources", "PDF • DOCX • TXT • HTML • CSV • PNG"); x += bw2 + g2
f2 = box(s, x, y, bw2, bh2, "Extract", "pdfminer/PyMuPDF • python-docx • BeautifulSoup"); x += bw2 + g2
f3 = box(s, x, y, bw2, bh2, "OCR (scans/PNG)", "pytesseract • OpenCV"); x += bw2 + g2
f4 = box(s, x, y, bw2, bh2, "To text", "merge regions")
for a,b in [(f1,f2),(f2,f3),(f3,f4)]: connect(s,a,b)
pipe = box(s, SAFE_L + 2.7, y + bh2 + 0.5, 4.2, bh2, "Scrub pipeline", "detect → explain → replace")
rend = box(s, SAFE_L, y + bh2 + 0.5, bw2, bh2, "Redacted output", "text replace • image masks + map")
connect(s, pipe, rend)
bullets(s, SAFE_L, SAFE_T + 2.2, SAFE_W, SAFE_H - 2.35, "Highlights", [
    "Kickoff priority: screenshots (PNG); PDFs use text layer first",
    "Both paths feed the same scrubber for consistent policy",
    "Image outputs can include a 'mask map' JSON for evidence",
], bullet_size=16)

s = title_only("Audit & governance (from kickoff)")
bullets(s, SAFE_L, SAFE_T, SAFE_W/2 - 0.05, SAFE_H, "What we log (examples)", [
    "Timestamp, corporate key, session open/close, logon/logoff",
    "Client, device, browser, (optional) location/MAC",
    "Original & scrubbed hashes; entities & confidences",
    "Highlight malicious prompts (search/export/update/delete)",
    "Receipts downloadable as JSON/JSONL",
], bullet_size=16)
bullets(s, SAFE_L + SAFE_W/2 + 0.05, SAFE_T, SAFE_W/2 - 0.05, SAFE_H, "De-scrub (role-gated)", [
    "Full or selective by entity IDs",
    "Mandatory justification (we implement even if 'nice-to-have')",
    "Every attempt logged and hash-chained",
    "Supports incident response & compliance checks",
], bullet_size=16)

s = title_only("API surface & examples")
bullets(s, SAFE_L, SAFE_T, SAFE_W/2 - 0.05, SAFE_H, "Endpoints", [
    "POST /scrub (text)",
    "POST /files/redact-text (file)",
    "POST /descrub (IDs + justification)",
    "GET /audit • GET /audit/jsonl • GET /docs (Swagger)",
], bullet_size=18)
code = (
    "curl -X POST /scrub -H 'Content-Type: application/json' "
    "-d '{\"text\":\"Email a@b.com\"}'\n"
    "→ { \"scrubbed\":\"Email <C3::EMAIL::…>\", \"receipt\":\"…\" }\n\n"
    "curl -X POST /descrub -H 'Content-Type: application/json' "
    "-d '{\"ids\":[\"…\"],\"justification\":\"fraud case\"}'\n"
    "→ { \"descrubbed\":\"Email a@b.com\" }"
)
tb = s.shapes.add_textbox(Inches(SAFE_L + SAFE_W/2 + 0.05), Inches(SAFE_T), Inches(SAFE_W/2 - 0.05), Inches(SAFE_H)).text_frame
tb.clear(); tb.word_wrap = True
p = tb.paragraphs[0]; p.text = code; p.font.name = "Courier New"; p.font.size = Pt(14); p.font.color.rgb = TEXT

s = title_only("UI preview (combined)")
col_w = (SAFE_W - 0.8) / 3.0; img_h = SAFE_H - 0.55; y0 = SAFE_T
for i, (pth, cap) in enumerate([(IMAGE1, "Interactive Scrubbing"), (IMAGE2, "Audit Dashboard"), (IMAGE3, "API Docs (Swagger)")]):
    x0 = SAFE_L + i*0.4 + col_w*i
    image_fit(s, pth, x0, y0, col_w, img_h, caption=cap)

s = title_only("KPI, milestones & deliverables (from kickoff)")
bullets(s, SAFE_L, SAFE_T, SAFE_W/2 - 0.05, SAFE_H, "Milestones & KPI", [
    "03/10 — M1: intermediary presentation + Q&A; KPI = % correctly scrubbed (recall-first)",
    "10/10 — M2 @ ING: pass/fail — all sensitive info scrubbed per parameters",
    "M2 deck: performance, strengths, weaknesses, improvements",
], bullet_size=16)
bullets(s, SAFE_L + SAFE_W/2 + 0.05, SAFE_T, SAFE_W/2 - 0.05, SAFE_H, "Deliverables", [
    "Python 3 modules; open-source allowed",
    "Unit tests for significant methods; precision/recall on golden set",
    "Performance (speed/efficiency) and documentation",
], bullet_size=16)

s = title_only("Tech primer + next steps & Q&A")
bullets(s, SAFE_L, SAFE_T, SAFE_W/2 - 0.05, SAFE_H, "Python tech (what & why)", [
    "FastAPI + Pydantic — REST + typed schemas; /docs via OpenAPI",
    "Regex + python-stdnum + Luhn — find + validate IBAN/card numbers",
    "phonenumbers + rapidfuzz — parse phones; context boosting",
    "pdfminer/PyMuPDF + python-docx + BeautifulSoup — extract before OCR",
    "pytesseract + OpenCV — OCR for screenshots/scans",
    "cryptography.Fernet + BLAKE3/SHA-256 — optional encrypted originals + tamper-evident chain",
    "pytest + ruff/black/mypy — tests & quality gates (golden set metrics)",
], bullet_size=16)
bullets(s, SAFE_L + SAFE_W/2 + 0.05, SAFE_T, SAFE_W/2 - 0.05, SAFE_H, "Next steps & decisions needed", [
    "Wire detectors + OCR writers; selective de-scrub with roles",
    "Nightly metrics with receipts (precision/recall per label)",
    "Confirm: label→C-level mapping; acceptable FP/FN; file mix; de-scrub governance; receipt retention",
], bullet_size=16)

prs.save(OUTFILE)
print(f"Written: {OUTFILE}")
