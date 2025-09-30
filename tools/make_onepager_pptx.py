from pptx import Presentation
from pptx.util import Inches, Pt

title = "SecurePrompt — MVP Demo (Week 1)"
bullets = [
    "Project & Objectives: Banking-grade prompt/file scrubbing with C2–C4 policies",
    "Milestones: MVP-3 demo (03 Oct), Final pass/fail (10 Oct)",
    "Architecture: detectors → scrub → policy → audit → de-scrub",
    "MVP Ladder: prompts → files (pdf/png OCR) → adaptive policies → API & metrics",
    "Demo: scrub prompt/pdf/png, flip C-level, show audit, selective de-scrub",
    "Metrics: precision/recall on golden set; OCR latency; throughput",
    "Risks: OCR quality, false positives, performance, audit integrity",
    "Next: classifier & dashboard; hardening",
]

prs = Presentation()
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = title

left, top, width, height = Inches(0.6), Inches(1.5), Inches(12), Inches(5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True
for i, b in enumerate(bullets):
    p = tf.add_paragraph() if i>0 else tf.paragraphs[0]
    p.text = "• " + b
    p.font.size = Pt(20)

prs.save("slides/OnePager.pptx")
print("Wrote slides/OnePager.pptx")
